"""
converters/pptx.py — Microsoft PowerPoint PPTX to Markdown converter.

Each slide is rendered as a separate page.  Text frames are converted to
Markdown (titles get a ``#`` prefix).  Images are extracted as base64-encoded
``PageImage`` objects.  Tables are rendered as GFM Markdown tables.
"""

from __future__ import annotations

import base64
import io
import logging
from pathlib import Path

from scinr.newton.converters.base import (
    BaseConverter,
    ConversionError,
    IntermediateDocument,
    PageImage,
)

logger = logging.getLogger(__name__)

# Numeric value for MSO_SHAPE_TYPE.PICTURE (fallback when enum is unavailable)
_MSO_PICTURE_TYPE = 13

_MEDIA_TYPE_MAP: dict[str, str] = {
    "JPEG": "image/jpeg",
    "PNG": "image/png",
    "GIF": "image/gif",
    "WEBP": "image/webp",
    "BMP": "image/bmp",
}


class PptxConverter(BaseConverter):
    """Convert ``.pptx`` files to the intermediate format.

    Each slide becomes a separate ``IntermediatePage``.  Text is extracted
    from all text frames (title shapes get a ``#`` prefix).  Images are
    base64-encoded and stored as ``PageImage`` objects.  Tables are
    converted to GFM Markdown tables.

    Requires ``python-pptx`` and optionally ``Pillow`` for media-type
    detection (falls back to ``image/png`` if Pillow is unavailable).
    """

    supported_extensions: frozenset[str] = frozenset({"pptx"})

    def convert(self, source: Path) -> IntermediateDocument:
        """Convert a PPTX file to the intermediate format.

        Parameters
        ----------
        source:
            Path to the ``.pptx`` file.

        Returns
        -------
        IntermediateDocument
            Multi-page document, one page per slide.

        Raises
        ------
        ConversionError
            If ``python-pptx`` is not installed or the file cannot be parsed.
        FileNotFoundError
            If *source* does not exist.
        """
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError as exc:
            raise ConversionError(
                "python-pptx is required for PPTX conversion. "
                "Install it with: uv add python-pptx"
            ) from exc

        if not source.exists():
            raise FileNotFoundError(f"File not found: {source}")

        try:
            prs = Presentation(str(source))
        except Exception as exc:
            raise ConversionError(f"Cannot open PPTX file {source}: {exc}") from exc

        try:
            pages = []
            for slide_idx, slide in enumerate(prs.slides):
                md_text = _extract_slide_markdown(slide)
                images = _extract_slide_images(slide)
                page = self._make_page(
                    index=slide_idx,
                    markdown=md_text,
                    images=images,
                )
                pages.append(page)
        except ConversionError:
            raise
        except Exception as exc:
            raise ConversionError(f"Failed to convert PPTX {source}: {exc}") from exc

        if not pages:
            logger.warning("%s has no slides — producing empty document.", source.name)
            pages.append(self._make_page(index=0, markdown=""))

        return IntermediateDocument(pages=pages)


# ---------------------------------------------------------------------------
# Private helpers
# ---------------------------------------------------------------------------


def _is_title_shape(shape) -> bool:  # type: ignore[no-untyped-def]
    """Return ``True`` if *shape* is a title placeholder.

    Parameters
    ----------
    shape:
        A ``python-pptx`` shape object.

    Returns
    -------
    bool
        Whether the shape should be treated as a slide title.
    """
    # Check by name convention first
    if "title" in shape.name.lower():
        return True
    # Check by placeholder type when available
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore[import-untyped]

        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                return True
    except (ImportError, AttributeError):
        pass
    return False


def _extract_slide_markdown(slide) -> str:  # type: ignore[no-untyped-def]
    """Extract and convert all text from a slide to Markdown.

    Parameters
    ----------
    slide:
        A ``python-pptx`` Slide object.

    Returns
    -------
    str
        GFM Markdown string with the slide's textual content.
    """
    # Sort shapes top-to-bottom, left-to-right
    sorted_shapes = sorted(
        slide.shapes, key=lambda s: (s.top or 0, s.left or 0)
    )

    parts: list[str] = []
    for shape in sorted_shapes:
        if shape.has_text_frame:
            paragraphs_text = "\n".join(
                para.text
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            )
            if not paragraphs_text.strip():
                continue
            if _is_title_shape(shape):
                parts.append(f"# {paragraphs_text.strip()}")
            else:
                parts.append(paragraphs_text.strip())
        elif shape.has_table:
            md_table = _pptx_table_to_markdown(shape.table)
            if md_table:
                parts.append(md_table)

    return "\n\n".join(p for p in parts if p)


def _extract_slide_images(slide) -> list[PageImage]:  # type: ignore[no-untyped-def]
    """Extract all images from a slide as ``PageImage`` objects.

    Parameters
    ----------
    slide:
        A ``python-pptx`` Slide object.

    Returns
    -------
    list[PageImage]
        List of base64-encoded images extracted from the slide.
    """
    images: list[PageImage] = []
    img_idx = 0

    for shape in slide.shapes:
        # Primary detection: check shape_type against MSO_SHAPE_TYPE.PICTURE
        is_picture = False
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]

            is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        except (ImportError, AttributeError):
            is_picture = shape.shape_type == _MSO_PICTURE_TYPE

        # Fallback: check for the presence of the image attribute
        if not is_picture:
            is_picture = hasattr(shape, "image")

        if not is_picture:
            continue

        try:
            image_blob: bytes = shape.image.blob
        except AttributeError:
            logger.debug(
                "Shape %r has no accessible image blob — skipping.", shape.name
            )
            continue

        media_type = _detect_media_type(image_blob)
        b64 = base64.b64encode(image_blob).decode("utf-8")
        images.append(
            PageImage(
                index=img_idx,
                base64=b64,
                media_type=media_type,
                description="",
            )
        )
        img_idx += 1

    return images


def _detect_media_type(blob: bytes) -> str:
    """Detect the MIME type of an image from its raw bytes.

    Uses Pillow for detection when available; falls back to
    ``"image/png"`` if Pillow is not installed or detection fails.

    Parameters
    ----------
    blob:
        Raw image bytes.

    Returns
    -------
    str
        MIME type string (e.g. ``"image/jpeg"``).
    """
    try:
        from PIL import Image  # type: ignore[import-untyped]

        img = Image.open(io.BytesIO(blob))
        fmt = img.format
        if fmt:
            return _MEDIA_TYPE_MAP.get(fmt.upper(), "image/png")
    except ImportError:
        logger.debug(
            "Pillow is not installed; defaulting to image/png for media-type detection. "
            "Install it with: uv add Pillow"
        )
    except Exception:
        pass
    return "image/png"


def _pptx_table_to_markdown(table) -> str:  # type: ignore[no-untyped-def]
    """Convert a ``python-pptx`` Table to a GFM Markdown table.

    Parameters
    ----------
    table:
        A ``python-pptx`` ``Table`` object.

    Returns
    -------
    str
        GFM Markdown table string, or empty string if the table has no rows.
    """
    all_rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        all_rows.append(cells)

    if not all_rows:
        return ""

    col_count = max(len(r) for r in all_rows)

    def _pad(row: list[str]) -> list[str]:
        padded = list(row)
        while len(padded) < col_count:
            padded.append("")
        return padded[:col_count]

    lines: list[str] = []
    header = _pad(all_rows[0])
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")
    for row in all_rows[1:]:
        lines.append("| " + " | ".join(_pad(row)) + " |")

    return "\n".join(lines)
